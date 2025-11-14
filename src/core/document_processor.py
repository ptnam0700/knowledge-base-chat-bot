"""
Document processing module for Thunderbolts.
Handles text extraction from various document formats with LangChain integration.
"""
from pathlib import Path
from typing import Optional, Union, List, Dict, Any
from collections import OrderedDict

try:
    import pdfplumber
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

# PowerPoint support (PPTX via python-pptx)
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# LangChain document loaders and text splitters (v0.2 compatible with fallbacks)
# Loaders moved to langchain_community; splitters moved to langchain_text_splitters
try:
    from langchain_community.document_loaders import (
        PyPDFLoader,
        TextLoader,
        Docx2txtLoader,
        UnstructuredFileLoader,
    )
    try:
        from langchain_community.document_loaders import PyMuPDFLoader  # Fast PDF loader
        PYMUPDF_LOADER_AVAILABLE = True
    except Exception:
        PyMuPDFLoader = None  # type: ignore
        PYMUPDF_LOADER_AVAILABLE = False
    try:
        from langchain_community.document_loaders import UnstructuredExcelLoader  # XLSX
        EXCEL_LOADER_AVAILABLE = True
    except Exception:
        UnstructuredExcelLoader = None  # type: ignore
        EXCEL_LOADER_AVAILABLE = False
    # PowerPoint-specific loader (PPT/PPTX)
    try:
        from langchain_community.document_loaders import UnstructuredPowerPointLoader  # type: ignore
        PPT_LOADER_AVAILABLE = True
    except Exception:
        UnstructuredPowerPointLoader = None  # type: ignore
        PPT_LOADER_AVAILABLE = False
    LOADER_AVAILABLE = True
except Exception:
    try:
        from langchain.document_loaders import (
            PyPDFLoader,
            TextLoader,
            Docx2txtLoader,
            UnstructuredFileLoader,  # type: ignore
        )
        try:
            from langchain.document_loaders import PyMuPDFLoader  # type: ignore
            PYMUPDF_LOADER_AVAILABLE = True
        except Exception:
            PyMuPDFLoader = None  # type: ignore
            PYMUPDF_LOADER_AVAILABLE = False
        try:
            from langchain.document_loaders import UnstructuredExcelLoader  # type: ignore
            EXCEL_LOADER_AVAILABLE = True
        except Exception:
            UnstructuredExcelLoader = None  # type: ignore
            EXCEL_LOADER_AVAILABLE = False
        # PowerPoint-specific loader (PPT/PPTX)
        try:
            from langchain.document_loaders import UnstructuredPowerPointLoader  # type: ignore
            PPT_LOADER_AVAILABLE = True
        except Exception:
            UnstructuredPowerPointLoader = None  # type: ignore
            PPT_LOADER_AVAILABLE = False
        LOADER_AVAILABLE = True
    except Exception:
        LOADER_AVAILABLE = False
        PYMUPDF_LOADER_AVAILABLE = False
        EXCEL_LOADER_AVAILABLE = False
        PPT_LOADER_AVAILABLE = False
        UnstructuredPowerPointLoader = None  # type: ignore

try:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
    SPLITTER_AVAILABLE = True
except Exception:
    try:
        from langchain.text_splitter import RecursiveCharacterTextSplitter  # type: ignore
        SPLITTER_AVAILABLE = True
    except Exception:
        SPLITTER_AVAILABLE = False

from .base_processor import BaseProcessor
from src.utils.exceptions import TextProcessingError, UnsupportedFormatError
from src.utils.performance import performance_context, measure_performance

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except Exception:
    PANDAS_AVAILABLE = False

# URL/text extraction dependencies (assumed present elsewhere in the project)
try:
    import requests
    from bs4 import BeautifulSoup
    WEB_AVAILABLE = True
except Exception:
    WEB_AVAILABLE = False


class DocumentProcessor(BaseProcessor):
    """Handles document text extraction operations."""
    
    def __init__(self, config: Optional[dict] = None):
        """
        Initialize the document processor.
        
        Args:
            config: Optional configuration dictionary
        """
        super().__init__(config)
        
        # Initialize LangChain text splitter if available (not used during extraction anymore)
        if SPLITTER_AVAILABLE:
            self.text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000,
                chunk_overlap=200,
                length_function=len,
                separators=["\n\n", "\n", " ", ""]
            )
        else:
            self.text_splitter = None
        # Extraction cache (in-memory, simple LRU)
        self.enable_extract_cache: bool = getattr(self.settings, "enable_extract_cache", True)
        self._extract_cache: "OrderedDict[str, str]" = OrderedDict()
        self._extract_cache_max_entries: int = 64
    
    @measure_performance
    def process(self, input_data: Union[str, Path], **kwargs) -> str:
        """
        Process document and extract text.
        
        Args:
            input_data: Path to document file or URL
            **kwargs: Additional processing parameters
            
        Returns:
            Extracted text content
            
        Raises:
            TextProcessingError: If processing fails
        """
        
        # Process as file
        if not self.validate_input(input_data):
            raise TextProcessingError(f"Invalid input: {input_data}")
        
        file_path = Path(input_data)
        
        if not self.validate_document_format(file_path):
            raise UnsupportedFormatError(f"Unsupported document format: {file_path.suffix}")
        
        self.log_processing_start("Document processing", str(file_path))
        
        try:
            # Cache lookup
            cache_key = None
            if self.enable_extract_cache and file_path.exists():
                try:
                    stat = file_path.stat()
                    cache_key = f"{str(file_path.resolve())}|{stat.st_size}|{int(stat.st_mtime)}"
                    if cache_key in self._extract_cache:
                        text_cached = self._extract_cache.pop(cache_key)
                        # move to end (recently used)
                        self._extract_cache[cache_key] = text_cached
                        return text_cached
                except Exception:
                    cache_key = None
            # Use LangChain loaders if available, else fallback to legacy
            if LOADER_AVAILABLE and kwargs.get('use_langchain', True):
                text = self._extract_text_with_langchain(file_path)
            else:
                text = self._extract_text_legacy(file_path)
            # Cache store
            if cache_key and self.enable_extract_cache and text:
                try:
                    self._extract_cache[cache_key] = text
                    if len(self._extract_cache) > self._extract_cache_max_entries:
                        # pop least recently used
                        self._extract_cache.popitem(last=False)
                except Exception:
                    pass
            self.log_processing_end("Document processing")
            return text
            
        except Exception as e:
            self.handle_error(e, "Document processing")

    def _extract_text_with_langchain(self, file_path: Path) -> str:
        """Extract text using LangChain document loaders."""
        suffix = file_path.suffix.lower()
        try:
            # Select loader per suffix
            if suffix == '.pdf':
                use_fast_loader = False
                try:
                    size_mb = file_path.stat().st_size / (1024 * 1024)
                    threshold = getattr(self.settings, 'pdf_large_file_size_mb', 15)
                    use_fast_loader = size_mb >= threshold
                except Exception:
                    use_fast_loader = False
                if PYMUPDF_LOADER_AVAILABLE and use_fast_loader:
                    loader = PyMuPDFLoader(str(file_path))  # type: ignore
                else:
                    loader = PyPDFLoader(str(file_path))
            elif suffix == '.docx':
                loader = Docx2txtLoader(str(file_path))
            elif suffix == '.txt':
                loader = TextLoader(str(file_path))
            elif suffix == '.xlsx':
                if EXCEL_LOADER_AVAILABLE:
                    loader = UnstructuredExcelLoader(str(file_path))  # type: ignore
                else:
                    # Fallback: read with pandas
                    return self.extract_text_from_xlsx(file_path)
            elif suffix in ('.ppt', '.pptx'):
                # Prefer dedicated PowerPoint loader if available
                if PPT_LOADER_AVAILABLE and UnstructuredPowerPointLoader is not None:
                    loader = UnstructuredPowerPointLoader(str(file_path))  # type: ignore
                else:
                    # Fallback to generic unstructured loader
                    loader = UnstructuredFileLoader(str(file_path))
            else:
                # Fallback to unstructured loader
                loader = UnstructuredFileLoader(str(file_path))

            with performance_context(f"loader_load_{suffix}"):
                documents = loader.load()

            # Combine all documents into raw text
            text = "\n\n".join([doc.page_content for doc in documents])
            return text

        except Exception as e:
            self.logger.warning(f"LangChain extraction failed for {file_path}: {e}")
            # Fallback to legacy method
            return self._extract_text_legacy(file_path)

    def _extract_text_legacy(self, file_path: Path) -> str:
        """Legacy text extraction method."""
        suffix = file_path.suffix.lower()
        
        if suffix == '.pdf':
            return self.extract_text_from_pdf(file_path)
        elif suffix == '.docx':
            return self.extract_text_from_docx(file_path)
        elif suffix == '.txt':
            return self.extract_text_from_txt(file_path)
        elif suffix == '.xlsx':
            return self.extract_text_from_xlsx(file_path)
        elif suffix in ('.ppt', '.pptx'):
            # Legacy PPTX extraction via python-pptx; .ppt best-effort (may fail)
            return self.extract_text_from_pptx(file_path)
        else:
            raise UnsupportedFormatError(f"Unsupported format: {suffix}")
    
    def extract_text_from_pdf(self, pdf_path: Path) -> str:
        """
        Extract text from PDF file.
        
        Args:
            pdf_path: Path to PDF file
            
        Returns:
            Extracted text
            
        Raises:
            TextProcessingError: If extraction fails
        """
        if not PDF_AVAILABLE:
            raise TextProcessingError("PDF processing not available. Please install pdfplumber.")
        
        try:
            text_content = []
            
            with pdfplumber.open(str(pdf_path)) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_content.append(page_text)
            
            return '\n\n'.join(text_content)
            
        except Exception as e:
            raise TextProcessingError(f"Failed to extract text from PDF: {e}")
    
    def extract_text_from_docx(self, docx_path: Path) -> str:
        """
        Extract text from DOCX file.
        
        Args:
            docx_path: Path to DOCX file
            
        Returns:
            Extracted text
            
        Raises:
            TextProcessingError: If extraction fails
        """
        if not DOCX_AVAILABLE:
            raise TextProcessingError("DOCX processing not available. Please install python-docx.")
        
        try:
            doc = Document(str(docx_path))
            text_content = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            return '\n\n'.join(text_content)
            
        except Exception as e:
            raise TextProcessingError(f"Failed to extract text from DOCX: {e}")

    def extract_text_from_pptx(self, ppt_path: Path) -> str:
        """
        Extract text from PPT/PPTX file using python-pptx as a legacy fallback.

        NOTE:
            - This supports .pptx reliably.
            - For .ppt (legacy binary format), python-pptx may fail; in that case
              a TextProcessingError is raised and upstream can fall back to LC loaders
              or treat the format as unsupported.
        """
        if not PPTX_AVAILABLE:
            raise TextProcessingError(
                "PPTX processing not available. Please install 'python-pptx'."
            )

        try:
            prs = Presentation(str(ppt_path))
            text_content: List[str] = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    # Most text-bearing shapes expose a .text attribute
                    if hasattr(shape, "text") and shape.text:
                        stripped = shape.text.strip()
                        if stripped:
                            text_content.append(stripped)

            return "\n\n".join(text_content)
        except Exception as e:
            raise TextProcessingError(f"Failed to extract text from PPT/PPTX: {e}")

    def extract_text_from_xlsx(self, xlsx_path: Path) -> str:
        """
        Extract text from XLSX file.
        Attempts to use pandas as a lightweight fallback when UnstructuredExcelLoader is unavailable.
        """
        try:
            if EXCEL_LOADER_AVAILABLE:
                # Prefer using the unstructured loader if available via langchain community
                with performance_context("unstructured_excel_load"):
                    loader = UnstructuredExcelLoader(str(xlsx_path))  # type: ignore
                    docs = loader.load()
                    return "\n\n".join([d.page_content for d in docs])
            if not PANDAS_AVAILABLE:
                raise TextProcessingError("XLSX processing not available. Install 'pandas' or 'unstructured'.")
            # Pandas fallback: concatenate all sheets and cells into text
            with performance_context("pandas_read_excel"):
                excel = pd.read_excel(str(xlsx_path), sheet_name=None, header=None)
            parts: List[str] = []
            for sheet_name, df in excel.items():
                try:
                    # Fill NaNs and convert to string; join rows by space and rows by newline
                    df_str = df.fillna("").astype(str)
                    rows = [" ".join(row) for row in df_str.values.tolist()]
                    sheet_text = f"[Sheet: {sheet_name}]\n" + "\n".join(rows)
                    parts.append(sheet_text)
                except Exception:
                    continue
            return "\n\n".join(parts)
        except Exception as e:
            raise TextProcessingError(f"Failed to extract text from XLSX: {e}")
    
    def extract_text_from_txt(self, txt_path: Path) -> str:
        """
        Extract text from TXT file.
        
        Args:
            txt_path: Path to TXT file
            
        Returns:
            Extracted text
            
        Raises:
            TextProcessingError: If extraction fails
        """
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            
            for encoding in encodings:
                try:
                    with open(txt_path, 'r', encoding=encoding) as file:
                        return file.read()
                except UnicodeDecodeError:
                    continue
            
            raise TextProcessingError("Could not decode text file with any supported encoding")
            
        except Exception as e:
            raise TextProcessingError(f"Failed to extract text from TXT: {e}")
    
    def extract_text_from_url(self, url: str) -> str:
        """
        Extract text from web URL.
        
        Args:
            url: Web URL to extract text from
            
        Returns:
            Extracted text
            
        Raises:
            TextProcessingError: If extraction fails
        """
        if not WEB_AVAILABLE:
            raise TextProcessingError("Web scraping not available. Please install beautifulsoup4 and requests.")
        
        try:
            # Fetch web page
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            }
            response = requests.get(url, headers=headers, timeout=30)
            response.raise_for_status()
            
            # Parse HTML and extract text
            soup = BeautifulSoup(response.content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Get text content
            text = soup.get_text()
            
            # Clean up text
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = ' '.join(chunk for chunk in chunks if chunk)
            
            return text
            
        except Exception as e:
            raise TextProcessingError(f"Failed to extract text from URL: {e}")

    def get_page_title(self, url: str) -> Optional[str]:
        """Fetch page title from URL if available."""
        if not WEB_AVAILABLE:
            return None
        try:
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            }
            response = requests.get(url, headers=headers, timeout=30)
            response.raise_for_status()
            soup = BeautifulSoup(response.content, 'html.parser')
            if soup.title and soup.title.string:
                return soup.title.string.strip()
            return None
        except Exception:
            return None
    
    def validate_document_format(self, file_path: Path) -> bool:
        """
        Validate if file format is supported for document processing.
        
        Args:
            file_path: Path to file
            
        Returns:
            True if format is supported
        """
        suffix = file_path.suffix.lower().lstrip('.')
        return suffix in self.settings.supported_document_formats
    
    def get_document_info(self, file_path: Path) -> dict:
        """
        Get document information.
        
        Args:
            file_path: Path to document file
            
        Returns:
            Dictionary with document information
        """
        try:
            info = {
                "file_size": file_path.stat().st_size,
                "format": file_path.suffix.lower(),
                "name": file_path.name
            }
            
            # Add format-specific info
            if file_path.suffix.lower() == '.pdf' and PDF_AVAILABLE:
                try:
                    with pdfplumber.open(str(file_path)) as pdf:
                        info["pages"] = len(pdf.pages)
                except:
                    pass
            
            return info
            
        except Exception as e:
            self.logger.warning(f"Failed to get document info: {e}")
            return {"error": str(e)}
